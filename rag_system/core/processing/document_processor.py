"""
Enhanced Document Processor with Multi-Format Support
Supports PDF, Word, PowerPoint, Excel, and text files
"""

import os
import io
from typing import Dict, List, Optional, Union
from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from rag_system.core.utils.logger import get_logger

logger = get_logger(__name__)

class EnhancedDocumentProcessor:
    """Process various document formats into text"""

    def __init__(self):
        self.supported_formats = {
            '.txt': self._process_text,
            '.md': self._process_text,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_doc_fallback,
            '.rtf': self._process_rtf,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_powerpoint,
            '.odt': self._process_odt
        }

        logger.info(f"Enhanced Document Processor initialized with {len(self.supported_formats)} supported formats")

    def process_file(self, file_path: Union[str, Path], file_content: Optional[bytes] = None) -> Dict:
        """
        Process a file and extract text content

        Args:
            file_path: Path to the file or filename
            file_content: Optional byte content of the file

        Returns:
            Dictionary with extracted content and metadata
        """
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()

        if file_extension not in self.supported_formats:
            logger.warning(f"Unsupported file format: {file_extension}")
            return {
                'content': '',
                'metadata': {'error': f'Unsupported format: {file_extension}'},
                'success': False
            }

        try:
            processor = self.supported_formats[file_extension]

            if file_content:
                # Process from byte content
                result = processor(file_content, is_bytes=True)
            else:
                # Process from file path
                result = processor(file_path)

            result['metadata'].update({
                'filename': file_path.name,
                'file_type': file_extension,
                'processor': 'enhanced_document_processor'
            })

            logger.info(f"Successfully processed {file_path.name} ({len(result['content'])} characters)")
            return result

        except Exception as e:
            logger.error(f"Error processing file {file_path.name}: {e}")
            return {
                'content': '',
                'metadata': {
                    'filename': file_path.name,
                    'file_type': file_extension,
                    'error': str(e)
                },
                'success': False
            }

    def _process_text(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process plain text files"""
        try:
            if is_bytes:
                content = source.decode('utf-8')
            else:
                with open(source, 'r', encoding='utf-8') as f:
                    content = f.read()

            return {
                'content': content,
                'metadata': {'format': 'text'},
                'success': True
            }
        except UnicodeDecodeError:
            # Try with different encoding
            if is_bytes:
                content = source.decode('latin-1')
            else:
                with open(source, 'r', encoding='latin-1') as f:
                    content = f.read()

            return {
                'content': content,
                'metadata': {'format': 'text', 'encoding': 'latin-1'},
                'success': True
            }

    def _process_pdf(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process PDF files"""
        if not HAS_PYPDF:
            return {
                'content': '',
                'metadata': {'error': 'pypdf not installed'},
                'success': False
            }

        try:
            if is_bytes:
                pdf_file = io.BytesIO(source)
            else:
                pdf_file = open(source, 'rb')

            reader = pypdf.PdfReader(pdf_file)
            content = []

            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    content.append(f"[Page {page_num + 1}]\n{text}")

            if not is_bytes:
                pdf_file.close()

            return {
                'content': '\n\n'.join(content),
                'metadata': {
                    'format': 'pdf',
                    'pages': len(reader.pages),
                    'title': reader.metadata.get('/Title', '') if reader.metadata else ''
                },
                'success': True
            }

        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'PDF processing failed: {e}'},
                'success': False
            }

    def _process_docx(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process Word DOCX files"""
        if not HAS_DOCX:
            return {
                'content': '',
                'metadata': {'error': 'python-docx not installed'},
                'success': False
            }

        try:
            if is_bytes:
                doc = DocxDocument(io.BytesIO(source))
            else:
                doc = DocxDocument(source)

            content = []

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)

            # Extract tables
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_content = [cell.text.strip() for cell in row.cells]
                    table_content.append(' | '.join(row_content))

                if table_content:
                    content.append('\n[Table]\n' + '\n'.join(table_content))

            return {
                'content': '\n\n'.join(content),
                'metadata': {
                    'format': 'docx',
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables)
                },
                'success': True
            }

        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'DOCX processing failed: {e}'},
                'success': False
            }

    def _process_doc_fallback(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Fallback for older DOC files"""
        return {
            'content': '',
            'metadata': {'error': 'Legacy DOC format not supported. Please convert to DOCX.'},
            'success': False
        }

    def _process_rtf(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Basic RTF processing"""
        try:
            if is_bytes:
                content = source.decode('utf-8')
            else:
                with open(source, 'r', encoding='utf-8') as f:
                    content = f.read()

            # Basic RTF cleaning (remove basic RTF tags)
            import re
            content = re.sub(r'\\[a-z]+\d*\s?', '', content)  # Remove RTF commands
            content = re.sub(r'[{}]', '', content)  # Remove braces
            content = content.replace('\\', '')

            return {
                'content': content,
                'metadata': {'format': 'rtf'},
                'success': True
            }
        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'RTF processing failed: {e}'},
                'success': False
            }

    def _process_csv(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process CSV files"""
        if not HAS_PANDAS:
            return {
                'content': '',
                'metadata': {'error': 'pandas not installed'},
                'success': False
            }

        try:
            if is_bytes:
                df = pd.read_csv(io.BytesIO(source))
            else:
                df = pd.read_csv(source)

            # Convert to readable text format
            content = []
            content.append(f"CSV Data with {len(df)} rows and {len(df.columns)} columns\n")
            content.append("Columns: " + ", ".join(df.columns.tolist()))
            content.append("\nData Preview:")
            content.append(df.head(10).to_string())

            if len(df) > 10:
                content.append(f"\n... and {len(df) - 10} more rows")

            return {
                'content': '\n'.join(content),
                'metadata': {
                    'format': 'csv',
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                },
                'success': True
            }

        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'CSV processing failed: {e}'},
                'success': False
            }

    def _process_excel(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process Excel files"""
        if not HAS_PANDAS:
            return {
                'content': '',
                'metadata': {'error': 'pandas not installed'},
                'success': False
            }

        try:
            if is_bytes:
                excel_file = pd.ExcelFile(io.BytesIO(source))
            else:
                excel_file = pd.ExcelFile(source)

            content = []
            sheet_info = []

            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)

                content.append(f"\n[Sheet: {sheet_name}]")
                content.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                content.append("Columns: " + ", ".join(df.columns.tolist()))
                content.append("\nData Preview:")
                content.append(df.head(5).to_string())

                sheet_info.append({
                    'name': sheet_name,
                    'rows': len(df),
                    'columns': len(df.columns)
                })

            return {
                'content': '\n'.join(content),
                'metadata': {
                    'format': 'excel',
                    'sheets': sheet_info,
                    'total_sheets': len(excel_file.sheet_names)
                },
                'success': True
            }

        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'Excel processing failed: {e}'},
                'success': False
            }

    def _process_powerpoint(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process PowerPoint files"""
        if not HAS_PPTX:
            return {
                'content': '',
                'metadata': {'error': 'python-pptx not installed'},
                'success': False
            }

        try:
            if is_bytes:
                prs = Presentation(io.BytesIO(source))
            else:
                prs = Presentation(source)

            content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"\n[Slide {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                content.extend(slide_content)

            return {
                'content': '\n'.join(content),
                'metadata': {
                    'format': 'powerpoint',
                    'slides': len(prs.slides)
                },
                'success': True
            }

        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'PowerPoint processing failed: {e}'},
                'success': False
            }

    def _process_odt(self, source: Union[Path, bytes], is_bytes: bool = False) -> Dict:
        """Process OpenDocument Text files"""
        try:
            if is_bytes:
                odt_file = zipfile.ZipFile(io.BytesIO(source))
            else:
                odt_file = zipfile.ZipFile(source)

            # Extract content.xml
            content_xml = odt_file.read('content.xml')
            root = ET.fromstring(content_xml)

            # Extract text content
            text_content = []
            for elem in root.iter():
                if elem.text:
                    text_content.append(elem.text.strip())

            odt_file.close()

            return {
                'content': '\n'.join(filter(None, text_content)),
                'metadata': {'format': 'odt'},
                'success': True
            }

        except Exception as e:
            return {
                'content': '',
                'metadata': {'error': f'ODT processing failed: {e}'},
                'success': False
            }

    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_formats.keys())

    def is_supported(self, file_extension: str) -> bool:
        """Check if a file format is supported"""
        return file_extension.lower() in self.supported_formats


# Global instance
document_processor = EnhancedDocumentProcessor()